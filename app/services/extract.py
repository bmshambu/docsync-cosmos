"""Text extraction + chunking.

Ported from the Cowork skill ``skills/rfp-data-prep/extract_text.py`` and made
importable: output directories and chunk settings are passed in rather than
hard-coded relative to a repo root. Logic (paragraph extraction per format,
word-overlap chunking with page/section refs) is unchanged.
"""

from __future__ import annotations

import json
import re
from pathlib import Path


def clean(text: str) -> str:
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def make_chunks(
    paragraphs: list[dict],
    doc_id: str,
    filename: str,
    chunk_size: int,
    overlap: int,
) -> list[dict]:
    """Word-level overlapping chunks carrying page + section refs for citation."""
    chunks: list[dict] = []
    chunk_id = 0
    buffer_words: list[str] = []
    buffer_meta: list[tuple[int, int | None, str]] = []

    for para in paragraphs:
        for w in para["text"].split():
            buffer_meta.append((len(buffer_words), para["page"], para["section"]))
            buffer_words.append(w)

        while len(buffer_words) >= chunk_size:
            chunk_words = buffer_words[:chunk_size]
            meta_slice = buffer_meta[:chunk_size]
            chunks.append(
                {
                    "chunk_id": f"{doc_id}_chunk_{chunk_id:04d}",
                    "doc_id": doc_id,
                    "filename": filename,
                    "page_start": meta_slice[0][1],
                    "page_end": meta_slice[-1][1],
                    "section": meta_slice[0][2],
                    "text": " ".join(chunk_words),
                }
            )
            chunk_id += 1
            buffer_words = buffer_words[chunk_size - overlap :]
            buffer_meta = buffer_meta[chunk_size - overlap :]

    if buffer_words:
        chunks.append(
            {
                "chunk_id": f"{doc_id}_chunk_{chunk_id:04d}",
                "doc_id": doc_id,
                "filename": filename,
                "page_start": buffer_meta[0][1] if buffer_meta else None,
                "page_end": buffer_meta[-1][1] if buffer_meta else None,
                "section": buffer_meta[0][2] if buffer_meta else "",
                "text": " ".join(buffer_words),
            }
        )

    return chunks


# ── Per-format extractors (stream-in-memory: bytes in, no disk) ───────────────
# All three libraries accept a binary file-like object / bytes, so blob content
# is parsed straight from memory — nothing is written to a local blob cache.

def extract_docx(data: bytes) -> list[dict]:
    import io
    from docx import Document

    doc = Document(io.BytesIO(data))
    paragraphs: list[dict] = []
    current_section = "Document Start"
    word_count = 0

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        if para.style and para.style.name and para.style.name.startswith("Heading"):
            current_section = text
        word_count += len(text.split())
        page_estimate = max(1, word_count // 300 + 1)
        paragraphs.append(
            {"text": text, "page": page_estimate, "section": current_section}
        )
    return paragraphs


def extract_pdf(data: bytes) -> list[dict]:
    import fitz  # PyMuPDF

    doc = fitz.open(stream=data, filetype="pdf")
    paragraphs: list[dict] = []
    current_section = "Document Start"

    for page_num, page in enumerate(doc, start=1):
        for block in page.get_text("blocks"):
            text = block[4].strip()
            if not text:
                continue
            if len(text) < 80 and text.isupper():
                current_section = text
            paragraphs.append(
                {"text": clean(text), "page": page_num, "section": current_section}
            )
    return paragraphs


def extract_pptx(data: bytes) -> list[dict]:
    import io
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    paragraphs: list[dict] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        current_section = f"Slide {slide_num}"
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                paragraphs.append(
                    {"text": text, "page": slide_num, "section": current_section}
                )
    return paragraphs


_EXTRACTORS = {
    ".docx": extract_docx,
    ".pdf": extract_pdf,
    ".pptx": extract_pptx,
    ".ppt": extract_pptx,
}


def extract_paragraphs(filename: str, data: bytes) -> list[dict] | None:
    """Dispatch to the right extractor by extension. None for unsupported."""
    ext = Path(filename).suffix.lower()
    extractor = _EXTRACTORS.get(ext)
    return extractor(data) if extractor else None


# ── Public API ────────────────────────────────────────────────────────────────

def _doc_id_of(filename: str) -> str:
    return Path(filename).stem.replace(" ", "_")


def process_bytes(
    filename: str,
    data: bytes,
    text_dir: Path,
    chunks_dir: Path,
    chunk_size: int,
    overlap: int,
    skip_existing: bool = True,
) -> dict | None:
    """Extract one document from IN-MEMORY bytes (no blob written to disk).

    Writes the small ``<doc>.txt`` / ``<doc>_chunks.json`` scratch (KB-scale)
    and returns a result dict that ALSO carries the chunk list, so the caller
    can persist to the store without a file read-back. None for unsupported.
    """
    doc_id = _doc_id_of(filename)
    ext = Path(filename).suffix.lower()
    if ext not in _EXTRACTORS:
        return None

    text_out = text_dir / f"{doc_id}.txt"
    chunks_out = chunks_dir / f"{doc_id}_chunks.json"

    if skip_existing and text_out.exists() and chunks_out.exists():
        chunks = json.loads(chunks_out.read_text(encoding="utf-8"))
        return {"doc_id": doc_id, "filename": filename, "skipped": True,
                "chunks": len(chunks), "chunks_data": chunks}

    paragraphs = extract_paragraphs(filename, data) or []
    text_dir.mkdir(parents=True, exist_ok=True)
    chunks_dir.mkdir(parents=True, exist_ok=True)
    text_out.write_text("\n\n".join(p["text"] for p in paragraphs), encoding="utf-8")

    chunks = make_chunks(paragraphs, doc_id, filename, chunk_size, overlap)
    chunks_out.write_text(
        json.dumps(chunks, indent=2, ensure_ascii=False), encoding="utf-8")

    return {"doc_id": doc_id, "filename": filename, "skipped": False,
            "paragraphs": len(paragraphs), "chunks": len(chunks),
            "chunks_data": chunks}


def process_file(
    path: Path,
    text_dir: Path,
    chunks_dir: Path,
    chunk_size: int,
    overlap: int,
    skip_existing: bool = True,
) -> dict | None:
    """Path-based wrapper (local files) — reads bytes, delegates to process_bytes."""
    ext = path.suffix.lower()
    if ext not in _EXTRACTORS:
        return None
    # skip check without reading the file
    doc_id = path.stem.replace(" ", "_")
    if skip_existing and (text_dir / f"{doc_id}.txt").exists() \
            and (chunks_dir / f"{doc_id}_chunks.json").exists():
        return process_bytes(path.name, b"", text_dir, chunks_dir, chunk_size, overlap, True)
    return process_bytes(path.name, path.read_bytes(), text_dir, chunks_dir,
                         chunk_size, overlap, skip_existing)


def extract_all_streamed(
    names: list[str],
    reader,                       # reader(name) -> bytes
    text_dir: Path,
    chunks_dir: Path,
    chunk_size: int,
    overlap: int,
    skip_existing: bool = True,
    on_progress=None,
) -> list[dict]:
    """Extract a batch by STREAMING each document's bytes from ``reader`` — the
    blob binary is never written to a local cache. ``on_progress(done, total,
    result)`` fires after each doc. Only the small text/chunk scratch is written.
    """
    text_dir.mkdir(parents=True, exist_ok=True)
    chunks_dir.mkdir(parents=True, exist_ok=True)

    results: list[dict] = []
    total = len(names)
    for i, name in enumerate(names, start=1):
        try:
            doc_id = _doc_id_of(name)
            # Skip-existing avoids even the download when scratch is present
            if skip_existing and (text_dir / f"{doc_id}.txt").exists() \
                    and (chunks_dir / f"{doc_id}_chunks.json").exists():
                r = process_bytes(name, b"", text_dir, chunks_dir, chunk_size, overlap, True)
            else:
                data = reader(name)
                r = process_bytes(name, data, text_dir, chunks_dir, chunk_size, overlap, False)
            if r is None:
                r = {"filename": name, "error": f"unsupported format {Path(name).suffix}"}
        except Exception as exc:  # noqa: BLE001 — log and continue
            r = {"filename": name, "error": str(exc)}
        results.append(r)
        if on_progress:
            on_progress(i, total, r)
    return results


def extract_all(
    doc_paths: list[Path],
    text_dir: Path,
    chunks_dir: Path,
    chunk_size: int,
    overlap: int,
    skip_existing: bool = True,
    on_progress=None,
) -> list[dict]:
    """Path-based batch (local FolderSource / tests). Streams from local files."""
    by_name = {p.name: p for p in doc_paths}
    return extract_all_streamed(
        [p.name for p in doc_paths],
        reader=lambda n: by_name[n].read_bytes(),
        text_dir=text_dir, chunks_dir=chunks_dir,
        chunk_size=chunk_size, overlap=overlap,
        skip_existing=skip_existing, on_progress=on_progress,
    )
